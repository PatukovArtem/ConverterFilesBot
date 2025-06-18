import logging
import os
import io
from io import BytesIO
from tempfile import mkdtemp
from pathlib import Path
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    CallbackQueryHandler,
    MessageHandler,
    filters,
    ContextTypes
)

try:
    from docx2pdf import convert as docx_to_pdf_convert

    DOCX_TO_PDF_SUPPORT = True
except ImportError:
    DOCX_TO_PDF_SUPPORT = False

try:
    from pdf2docx import Converter as PdfToDocxConverter

    PDF_TO_DOCX_SUPPORT = True
except ImportError:
    PDF_TO_DOCX_SUPPORT = False

try:
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

from PIL import Image, ImageOps


logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

user_states = {}

keyboard_main = [
    [InlineKeyboardButton("🖼 Картинки", callback_data="images")],
    [InlineKeyboardButton("📁 Файлы", callback_data="files")]
]
markup_main = InlineKeyboardMarkup(keyboard_main)

keyboard_images = [
    [InlineKeyboardButton("PNG → JPG", callback_data="png_to_jpg"),
     InlineKeyboardButton("JPG → PNG", callback_data="jpg_to_png")],
    [InlineKeyboardButton("WEBP → JPG", callback_data="webp_to_jpg"),
     InlineKeyboardButton("JPG → WEBP", callback_data="jpg_to_webp")],
    [InlineKeyboardButton("Черно-белое", callback_data="to_grayscale")],
    [InlineKeyboardButton("🔙 Назад", callback_data="back_main")]
]
markup_images = InlineKeyboardMarkup(keyboard_images)

keyboard_files = []
if DOCX_TO_PDF_SUPPORT:
    keyboard_files.append([InlineKeyboardButton("DOCX → PDF", callback_data="docx_to_pdf")])
if PDF_TO_DOCX_SUPPORT:
    keyboard_files.append([InlineKeyboardButton("PDF → DOCX", callback_data="pdf_to_docx")])
if PPTX_SUPPORT:
    keyboard_files.append([InlineKeyboardButton("PPTX → PDF", callback_data="pptx_to_pdf")])
keyboard_files.append([InlineKeyboardButton("🔙 Назад", callback_data="back_main")])
markup_files = InlineKeyboardMarkup(keyboard_files)


async def convert_docx_to_pdf(input_path: str, output_path: str) -> None:
    """Конвертирует DOCX в PDF (кое-как)"""
    docx_to_pdf_convert(input_path, output_path)


async def convert_pdf_to_docx(input_path: str, output_path: str) -> None:
    """Конвертирует PDF в DOCX (тоже кое-как)"""
    converter = PdfToDocxConverter(input_path)
    converter.convert(output_path, start=0, end=None)
    converter.close()


async def convert_pptx_to_pdf(input_pptx: bytes) -> BytesIO:
    """Конвертация PPTX в PDF (здесь вообще молчу)"""
    prs = Presentation(io.BytesIO(input_pptx))
    pdf_buffer = BytesIO()
    c = canvas.Canvas(pdf_buffer, pagesize=letter)

    y_start = 700
    for slide in prs.slides:
        c.setFont("Helvetica", 12)
        y_position = y_start

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.replace('\r', '').replace('\x0b', '\n')
                for line in text.split('\n'):
                    if line.strip():
                        c.drawString(100, y_position, line.strip())
                        y_position -= 20

        c.showPage()

    c.save()
    pdf_buffer.seek(0)
    return pdf_buffer


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

    await update.message.reply_text(
        "📁 Главное меню \n"
        "Выберите категорию конвертации:",
        reply_markup=markup_main
    )


async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

    query = update.callback_query
    await query.answer()

    user_id = query.from_user.id
    chat_id = query.message.chat_id
    message_id = query.message.message_id

    if query.data == "images":
        user_states[user_id] = {"menu": "images"}
        await context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text="🖼 Выберите тип конвертации картинок:",
            reply_markup=markup_images
        )

    elif query.data == "files":
        user_states[user_id] = {"menu": "files"}
        await context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text="📁 Выберите тип конвертации файлов:",
            reply_markup=markup_files
        )

    elif query.data == "back_main":
        if user_id in user_states:
            del user_states[user_id]
        await context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text="📁 Главное меню \n"
                 "Выберите категорию конвертации:",
            reply_markup=markup_main
        )

    elif query.data in ["png_to_jpg", "jpg_to_png", "webp_to_jpg", "jpg_to_webp", "to_grayscale"]:
        user_states[user_id] = {"menu": "images", "conversion": query.data}
        await context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text=f"🖼 Вы выбрали: {query.data.replace('_', ' ').title().replace('To', '→')}\n"
                 f"Теперь отправьте картинку для обработки",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Назад", callback_data="images")]])
        )

    elif query.data in ["docx_to_pdf", "pdf_to_docx", "pptx_to_pdf"]:
        user_states[user_id] = {"menu": "files", "conversion": query.data}
        await context.bot.edit_message_text(
            chat_id=chat_id,
            message_id=message_id,
            text=f"📄 Вы выбрали: {query.data.replace('_', ' ').title().replace('To', '→')}\n"
                 f"Отправьте файл для конвертации",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Назад", callback_data="files")]])
        )


async def handle_images(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

    user_id = update.message.from_user.id

    if user_id not in user_states or user_states[user_id].get("menu") != "images":
        await update.message.reply_text(
            "Сначала выберите режим конвертации через меню",
            reply_markup=markup_main
        )
        return

    mode = user_states[user_id].get("conversion")
    photo = update.message.photo[-1] if update.message.photo else None
    document = update.message.document if update.message.document else None

    if not photo and not document:
        await update.message.reply_text("Пожалуйста, отправьте жалкую картинку")
        return

    try:
        file = await (photo or document).get_file()
        file_bytes = BytesIO()
        await file.download_to_memory(out=file_bytes)
        file_bytes.seek(0)
        img = Image.open(file_bytes)

        output_bytes = BytesIO()

        if mode == "png_to_jpg":
            if img.format != "PNG":
                await update.message.reply_text("Отправьте PNG картинку", reply_markup=markup_images)
                return
            img = img.convert('RGB')
            output_bytes.name = "converted.jpg"
            img.save(output_bytes, format="JPEG", quality=95)

        elif mode == "jpg_to_png":
            if img.format not in ["JPEG", "JPG"]:
                await update.message.reply_text("Отправьте JPG картинку", reply_markup=markup_images)
                return
            output_bytes.name = "converted.png"
            img.save(output_bytes, format="PNG")

        elif mode == "webp_to_jpg":
            if img.format != "WEBP":
                await update.message.reply_text("Отправьте WEBP картинку", reply_markup=markup_images)
                return
            img = img.convert('RGB')
            output_bytes.name = "converted.jpg"
            img.save(output_bytes, format="JPEG", quality=95)

        elif mode == "jpg_to_webp":
            if img.format not in ["JPEG", "JPG"]:
                await update.message.reply_text("Отправьте JPG картинку", reply_markup=markup_images)
                return
            output_bytes.name = "converted.webp"
            img.save(output_bytes, format="WEBP", quality=90)

        elif mode == "to_grayscale":
            output_bytes.name = "grayscale.jpg"
            ImageOps.grayscale(img).save(output_bytes, format="JPEG")

        output_bytes.seek(0)

        await update.message.reply_document(
            document=output_bytes,
            caption="Готово!😉😉😉"
        )

        await update.message.reply_text(
            "🖼 Выберите тип конвертации картинки:",
            reply_markup=markup_images
        )

    except Exception as e:
        logger.error(f"Чет пошло не так 😱😓 при обработки картинки: {e}")
        await update.message.reply_text(
            "Одни ошибки при обработке картинок 😱😓",
            reply_markup=markup_images
        )


async def handle_files(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

    user_id = update.message.from_user.id

    if user_id not in user_states or user_states[user_id].get("menu") != "files":
        await update.message.reply_text(
            "Сначала выберите режим конвертации через меню",
            reply_markup=markup_main
        )
        return

    mode = user_states[user_id].get("conversion")
    document = update.message.document

    if not document:
        await update.message.reply_text("Пожалуйста, отправьте файл")
        return

    temp_dir = None
    try:

        file = await document.get_file()
        file_bytes = BytesIO()
        await file.download_to_memory(out=file_bytes)
        input_data = file_bytes.getvalue()
        file_ext = Path(document.file_name).suffix.lower() if document.file_name else ""


        temp_dir = mkdtemp()
        input_path = os.path.join(temp_dir, f"input{file_ext}")
        with open(input_path, 'wb') as f:
            f.write(input_data)


        if mode == "docx_to_pdf" and DOCX_TO_PDF_SUPPORT:
            if not file_ext == ".docx":
                raise ValueError("Требуется файл .docx 😱😓")

            output_path = os.path.join(temp_dir, "output.pdf")
            await convert_docx_to_pdf(input_path, output_path)

            with open(output_path, 'rb') as f:
                await update.message.reply_document(
                    document=f,
                    filename=document.file_name.replace('.docx', '.pdf'),
                    caption="✅ PDF готов!😉😉😉"
                )

        elif mode == "pdf_to_docx" and PDF_TO_DOCX_SUPPORT:
            if not file_ext == ".pdf":
                raise ValueError("Требуется файл .pdf 😱😓")

            output_path = os.path.join(temp_dir, "output.docx")
            await convert_pdf_to_docx(input_path, output_path)

            with open(output_path, 'rb') as f:
                await update.message.reply_document(
                    document=f,
                    filename=document.file_name.replace('.pdf', '.docx'),
                    caption="✅ DOCX готов!😉😉😉"
                )

        elif mode == "pptx_to_pdf" and PPTX_SUPPORT:
            if not file_ext == ".pptx":
                raise ValueError("Требуется файл .pptx 😱😓")

            pdf_buffer = await convert_pptx_to_pdf(input_data)
            await update.message.reply_document(
                document=pdf_buffer,
                filename=document.file_name.replace('.pptx', '.pdf'),
                caption="✅ PDF готов!😉😉😉"
            )

        else:
            await update.message.reply_text(
                "Выбранный тип конвертации временно недоступен 😱😓",
                reply_markup=markup_files
            )
            return


        await update.message.reply_text(
            "📁 Выберите тип конвертации файлов:",
            reply_markup=markup_files
        )

    except ValueError as e:
        await update.message.reply_text(
            f"Ошибка: {str(e)}",
            reply_markup=markup_files
        )
    except Exception as e:
        logger.error(f"Ошибка 😱😓 конвертации файла: {e}")
        await update.message.reply_text(
            "Одни ошибки при обработке файлов 😱😓",
            reply_markup=markup_files
        )
    finally:
        if temp_dir and os.path.exists(temp_dir):
            for f in os.listdir(temp_dir):
                os.unlink(os.path.join(temp_dir, f))
            os.rmdir(temp_dir)


def main() -> None:

    application = ApplicationBuilder().token("7895923677:AAE598bRHaII-pdUyWupicxzl5ilqzZoZJo").build()


    application.add_handler(CommandHandler("start", start))
    application.add_handler(CallbackQueryHandler(button_handler))
    application.add_handler(MessageHandler(filters.PHOTO | filters.Document.IMAGE, handle_images))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_files))


    application.run_polling()


if __name__ == "__main__":
    main()